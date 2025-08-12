# -*- coding: utf-8 -*-
"""
搜索服务
"""
import os
import re
import json
from typing import List, Dict, Any, Optional
from pathlib import Path

# 导入智能编码检测工具
from app.utils.encoding_detector import EncodingDetector

# 导入文档解析库
try:
    import pandas as pd
except ImportError:
    pd = None

try:
    import PyPDF2
    import pdfplumber
except ImportError:
    PyPDF2 = None
    pdfplumber = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

# Windows COM接口用于处理老版本.doc文件
try:
    import win32com.client
    HAS_WIN32COM = True
except ImportError:
    win32com = None
    HAS_WIN32COM = False

# docx2txt作为备用文档提取工具
try:
    import docx2txt
    HAS_DOCX2TXT = True
except ImportError:
    docx2txt = None
    HAS_DOCX2TXT = False

# olefile用于直接解析OLE文档结构
try:
    import olefile
    HAS_OLEFILE = True
except ImportError:
    olefile = None
    HAS_OLEFILE = False

# msoffcrypto用于处理加密Office文档
try:
    import msoffcrypto
    HAS_MSOFFCRYPTO = True
except ImportError:
    msoffcrypto = None
    HAS_MSOFFCRYPTO = False

try:
    import markdown
except ImportError:
    markdown = None

class SearchService:
    """文件内容搜索服务"""
    
    def __init__(self):
        self.supported_extensions = {
            # 文本文件
            '.txt', '.md', '.csv', '.json', '.py', '.js', '.html', '.xml', '.yml', '.yaml',
            # 配置文件
            '.conf', '.config', '.cfg', '.ini', '.properties', '.env',
            # Office文档
            '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx',
            # 其他文档
            '.rtf', '.odt', '.ods', '.odp'
        }
    
    def search_in_text(self, text: str, keyword: str, quick_mode: bool = False) -> List[Dict[str, Any]]:
        """在文本内容中搜索关键词（公共接口）"""
        return self._search_in_text(text, keyword, quick_mode)
    
    def search_in_file(self, file_path: str, keyword: str, quick_mode: bool = False) -> List[Dict[str, Any]]:
        """在文件中搜索关键词"""
        try:
            if not os.path.exists(file_path):
                return []
            
            file_ext = Path(file_path).suffix.lower()
            if file_ext not in self.supported_extensions:
                return []
            
            # 检查文件大小
            file_size = os.path.getsize(file_path)
            
            # 快速模式：跳过大型复杂文件，但允许小型文档
            if quick_mode:
                if file_ext == '.pptx':  # PowerPoint始终跳过
                    return []
                if file_ext == '.pdf' and file_size > 1024 * 1024:  # PDF大于1MB时跳过
                    return []
                if file_ext in {'.docx', '.xlsx'} and file_size > 1024 * 1024:  # Office文档大于1MB时跳过
                    return []
            if file_size > 2 * 1024 * 1024:  # 大于2MB的文件
                if quick_mode:
                    return []
                # 对大文件只读取前10000字符
                content = self.extract_file_content_partial(file_path, max_chars=10000)
            else:
                # 读取文件内容
                content = self.extract_file_content(file_path)
            
            if not content:
                return []
            
            return self._search_in_text(content, keyword, quick_mode)
            
        except Exception as e:
            print(f"搜索文件 {file_path} 失败: {str(e)}")
            return []
    
    def _search_in_text(self, text: str, keyword: str, quick_mode: bool = False) -> List[Dict[str, Any]]:
        """在文本中搜索关键词"""
        matches = []
        lines = text.split('\n')
        
        # 支持正则表达式搜索
        try:
            pattern = re.compile(keyword, re.IGNORECASE)
        except re.error:
            # 如果不是有效的正则表达式，使用普通字符串搜索
            pattern = re.compile(re.escape(keyword), re.IGNORECASE)
        
        for line_num, line in enumerate(lines, 1):
            if pattern.search(line):
                # 快速模式下限制匹配数量
                if quick_mode and len(matches) >= 3:
                    break
                    
                # 提取匹配的上下文
                context_start = max(0, line_num - 3)
                context_end = min(len(lines), line_num + 2)
                context_lines = lines[context_start:context_end]
                
                # 高亮匹配的关键词
                highlighted_line = pattern.sub(
                    lambda m: f"<mark>{m.group()}</mark>", 
                    line
                )
                
                matches.append({
                    "line_number": line_num,
                    "text": highlighted_line,
                    "context": "\n".join(context_lines) if not quick_mode else "",
                    "match_position": line.lower().find(keyword.lower())
                })
        
        return matches
    
    def extract_file_content(self, file_path: str) -> Optional[str]:
        """提取文件内容"""
        try:
            if not os.path.exists(file_path):
                return None
            
            file_ext = Path(file_path).suffix.lower()
            
            # 文本文件直接读取 - 使用智能编码检测
            if file_ext in {'.txt', '.py', '.js', '.html', '.xml', '.yml', '.yaml', '.csv', '.rtf', 
                           '.conf', '.config', '.cfg', '.ini', '.properties', '.env'}:
                content, error = EncodingDetector.read_file_with_encoding(file_path)
                if content is not None:
                    return content
                else:
                    # 如果智能检测失败，使用原来的方法作为后备
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        return f.read()
            
            # Markdown文件处理 - 使用智能编码检测
            elif file_ext == '.md':
                content, error = EncodingDetector.read_file_with_encoding(file_path)
                if content is None:
                    # 如果智能检测失败，使用原来的方法作为后备
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    # 如果有markdown库，可以转换为纯文本
                    if markdown:
                        try:
                            # 先转换为HTML，然后提取纯文本
                            html = markdown.markdown(content)
                            # 简单地移除HTML标签
                            import re
                            text = re.sub(r'<[^>]+>', '', html)
                            return content + "\n\n[纯文本内容]\n" + text
                        except Exception:
                            pass
                    return content
            
            # JSON文件格式化读取
            elif file_ext == '.json':
                with open(file_path, 'r', encoding='utf-8') as f:
                    try:
                        data = json.load(f)
                        return json.dumps(data, indent=2, ensure_ascii=False)
                    except json.JSONDecodeError:
                        # 如果JSON格式错误，按文本读取
                        f.seek(0)
                        return f.read()
            
            # PDF文件处理
            elif file_ext == '.pdf':
                return self._extract_pdf_content(file_path)
            
            # Word文档处理
            elif file_ext in {'.doc', '.docx'}:
                return self._extract_word_content(file_path)
            
            # Excel文件处理
            elif file_ext in {'.xlsx', '.xls'}:
                return self._extract_excel_content(file_path)
            
            # PowerPoint文件处理
            elif file_ext in {'.ppt', '.pptx'}:
                return self._extract_powerpoint_content(file_path)
            
            return None
                    
        except Exception as e:
            print(f"提取文件内容失败 {file_path}: {str(e)}")
            return None
    
    def _extract_pdf_content(self, file_path: str) -> Optional[str]:
        """提取PDF文件内容"""
        text_content = []
        
        # 方法1: 使用pdfplumber (更准确)
        if pdfplumber:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        text = page.extract_text()
                        if text:
                            text_content.append(f"[页面 {page_num}]\n{text}\n")
                    
                    if text_content:
                        return "\n".join(text_content)
            except Exception as e:
                print(f"使用pdfplumber提取PDF失败: {e}")
        
        # 方法2: 使用PyPDF2 (备用)
        if PyPDF2 and not text_content:
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page_num, page in enumerate(pdf_reader.pages, 1):
                        text = page.extract_text()
                        if text:
                            text_content.append(f"[页面 {page_num}]\n{text}\n")
                    
                    if text_content:
                        return "\n".join(text_content)
            except Exception as e:
                print(f"使用PyPDF2提取PDF失败: {e}")
        
        return "PDF文件内容提取失败" if not text_content else None
    
    def _extract_word_content(self, file_path: str, retry_count: int = 0) -> Optional[str]:
        """提取Word文档内容"""
        
        max_retries = 2  # 最多重试2次
        method_results = {}  # 记录每个方法的尝试结果
        
        # 创建详细的诊断信息
        diagnostic_info = {
            'file_path': file_path,
            'file_exists': os.path.exists(file_path),
            'file_size': os.path.getsize(file_path) if os.path.exists(file_path) else 0,
            'file_extension': os.path.splitext(file_path)[1].lower(),
            'attempts': [],
            'final_result': None
        }
        
        print(f"\n=== .doc文件处理诊断报告 ===")
        print(f"文件路径: {diagnostic_info['file_path']}")
        print(f"文件存在: {diagnostic_info['file_exists']}")
        print(f"文件大小: {diagnostic_info['file_size']} 字节")
        print(f"文件扩展名: {diagnostic_info['file_extension']}")
        
        # 首先检查文件扩展名，选择合适的处理方法
        if file_path.lower().endswith('.doc'):
            # 老版本.doc文件，优先使用LibreOffice系统工具
            print(f"\n--- 尝试方法1: LibreOffice系统工具 ---")
            try:
                result = self._extract_doc_with_system_tools(file_path)
                if result:
                    diagnostic_info['attempts'].append({
                        'method': 'LibreOffice',
                        'status': 'success',
                        'result_length': len(result),
                        'details': 'LibreOffice成功提取内容'
                    })
                    diagnostic_info['final_result'] = 'LibreOffice'
                    print(f"[OK] LibreOffice处理成功，内容长度: {len(result)}")
                    return result
                else:
                    method_results['LibreOffice'] = '返回空内容'
                    diagnostic_info['attempts'].append({
                        'method': 'LibreOffice', 
                        'status': 'failed',
                        'error': '返回空内容',
                        'details': 'LibreOffice执行成功但未提取到内容'
                    })
                    print(f"[FAIL] LibreOffice处理失败: 返回空内容")
            except Exception as e:
                method_results['LibreOffice'] = str(e)
                diagnostic_info['attempts'].append({
                    'method': 'LibreOffice',
                    'status': 'error', 
                    'error': str(e),
                    'details': f'LibreOffice执行异常: {str(e)[:100]}'
                })
                print(f"[ERROR] LibreOffice处理异常: {str(e)[:100]}")
            
            # 备用方案：COM接口
            print(f"\n--- 尝试方法2: COM接口 ---")
            if HAS_WIN32COM:
                try:
                    result = self._extract_doc_content_with_com(file_path)
                    if result:
                        diagnostic_info['attempts'].append({
                            'method': 'COM接口',
                            'status': 'success',
                            'result_length': len(result),
                            'details': 'COM接口成功提取内容'
                        })
                        diagnostic_info['final_result'] = 'COM接口'
                        print(f"[OK] COM接口处理成功，内容长度: {len(result)}")
                        return result
                    else:
                        method_results['COM接口'] = '返回空内容'
                        diagnostic_info['attempts'].append({
                            'method': 'COM接口',
                            'status': 'failed',
                            'error': '返回空内容',
                            'details': 'COM接口执行成功但未提取到内容'
                        })
                        print(f"[FAIL] COM接口处理失败: 返回空内容")
                except Exception as e:
                    method_results['COM接口'] = str(e)
                    diagnostic_info['attempts'].append({
                        'method': 'COM接口',
                        'status': 'error',
                        'error': str(e),
                        'details': f'COM接口执行异常: {str(e)[:100]}'
                    })
                    print(f"[ERROR] COM接口处理异常: {str(e)[:100]}")
            else:
                method_results['COM接口'] = '未安装win32com'
                diagnostic_info['attempts'].append({
                    'method': 'COM接口',
                    'status': 'unavailable',
                    'error': '未安装win32com',
                    'details': 'win32com库未安装，无法使用COM接口'
                })
                print(f"[SKIP] COM接口不可用: 未安装win32com")
            
            # 尝试使用docx2txt作为备用方法
            print(f"\n--- 尝试方法3: docx2txt ---")
            if HAS_DOCX2TXT:
                try:
                    result = self._extract_doc_with_docx2txt(file_path)
                    if result:
                        diagnostic_info['attempts'].append({
                            'method': 'docx2txt',
                            'status': 'success',
                            'result_length': len(result),
                            'details': 'docx2txt成功提取内容'
                        })
                        diagnostic_info['final_result'] = 'docx2txt'
                        print(f"[OK] docx2txt处理成功，内容长度: {len(result)}")
                        return result
                    else:
                        method_results['docx2txt'] = '返回空内容'
                        diagnostic_info['attempts'].append({
                            'method': 'docx2txt',
                            'status': 'failed',
                            'error': '返回空内容',
                            'details': 'docx2txt执行成功但未提取到内容'
                        })
                        print(f"[FAIL] docx2txt处理失败: 返回空内容")
                except Exception as e:
                    method_results['docx2txt'] = str(e)
                    diagnostic_info['attempts'].append({
                        'method': 'docx2txt',
                        'status': 'error',
                        'error': str(e),
                        'details': f'docx2txt执行异常: {str(e)[:100]}'
                    })
                    print(f"[ERROR] docx2txt处理异常: {str(e)[:100]}")
            else:
                method_results['docx2txt'] = '未安装docx2txt'
                diagnostic_info['attempts'].append({
                    'method': 'docx2txt',
                    'status': 'unavailable',
                    'error': '未安装docx2txt',
                    'details': 'docx2txt库未安装'
                })
                print(f"[SKIP] docx2txt不可用: 未安装docx2txt")
            
            # 尝试使用olefile直接解析OLE结构
            print(f"\n--- 尝试方法4: olefile ---")
            if HAS_OLEFILE:
                try:
                    result = self._extract_doc_with_olefile(file_path)
                    if result:
                        diagnostic_info['attempts'].append({
                            'method': 'olefile',
                            'status': 'success',
                            'result_length': len(result),
                            'details': 'olefile成功解析OLE结构并提取内容'
                        })
                        diagnostic_info['final_result'] = 'olefile'
                        print(f"[OK] olefile处理成功，内容长度: {len(result)}")
                        return result
                    else:
                        method_results['olefile'] = '返回空内容'
                        diagnostic_info['attempts'].append({
                            'method': 'olefile',
                            'status': 'failed',
                            'error': '返回空内容',
                            'details': 'olefile执行成功但未提取到内容'
                        })
                        print(f"[FAIL] olefile处理失败: 返回空内容")
                except Exception as e:
                    method_results['olefile'] = str(e)
                    diagnostic_info['attempts'].append({
                        'method': 'olefile',
                        'status': 'error',
                        'error': str(e),
                        'details': f'olefile执行异常: {str(e)[:100]}'
                    })
                    print(f"[ERROR] olefile处理异常: {str(e)[:100]}")
            else:
                method_results['olefile'] = '未安装olefile'
                diagnostic_info['attempts'].append({
                    'method': 'olefile',
                    'status': 'unavailable',
                    'error': '未安装olefile',
                    'details': 'olefile库未安装'
                })
                print(f"[SKIP] olefile不可用: 未安装olefile")
            
            # 智能重试机制和最终错误报告
            if retry_count < max_retries:
                import time
                sleep_time = 2  # 等待2秒后重试，避免资源冲突
                print(f"\n--- 智能重试机制 ---")
                print(f"所有方法都失败了，{sleep_time}秒后进行第{retry_count + 1}次重试...")
                diagnostic_info['attempts'].append({
                    'method': '重试机制',
                    'status': 'retry',
                    'details': f'第{retry_count + 1}次重试，等待{sleep_time}秒'
                })
                time.sleep(sleep_time)
                return self._extract_word_content(file_path, retry_count + 1)
            
            # 生成详细的错误诊断报告
            print(f"\n=== 最终诊断报告 ===")
            print(f"处理结果: 失败")
            print(f"尝试次数: {len(diagnostic_info['attempts'])}")
            print(f"重试次数: {retry_count}")
            
            # 生成详细的错误报告
            error_details = []
            for method, error in method_results.items():
                error_details.append(f"{method}: {error}")
            
            detailed_error = (f"老版本.doc文件处理失败，详细诊断：\n" + 
                            f"=== .doc文件处理诊断报告 ===\n" + 
                            "\n".join(error_details) + f"\n重试次数: {retry_count}")
            
            # 根据错误类型提供建议
            suggestions = []
            if any('未安装' in str(error) for error in method_results.values()):
                suggestions.append("建议安装缺失的处理库")
            if 'LibreOffice' in method_results and '返回空内容' in method_results.get('LibreOffice', ''):
                suggestions.append("LibreOffice可能需要重启或检查安装")
            if diagnostic_info['file_size'] < 1024:
                suggestions.append("文件过小，可能是空文件或损坏文件")
            
            if suggestions:
                detailed_error += f"\n\n建议解决方案：\n" + "\n".join(f"- {s}" for s in suggestions)
            
            print(f"错误报告已生成: {len(detailed_error)} 字符")
            # 这里不抛出异常，而是继续尝试python-docx作为最后的fallback
        
        # 新版本.docx文件或COM接口失败的.doc文件
        if not DocxDocument:
            if file_path.lower().endswith('.doc'):
                return "老版本.doc文件需要Microsoft Word或win32com支持，请安装相关组件"
            else:
                return "需要安装python-docx库才能读取Word文档"
        
        try:
            print("[DEBUG] 尝试python-docx提取...")
            doc = DocxDocument(file_path)
            paragraphs = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # 提取表格内容
            tables_content = []
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    table_text.append(" | ".join(row_text))
                tables_content.append("\n".join(table_text))
            
            content = "\n".join(paragraphs)
            if tables_content:
                content += "\n\n[表格内容]\n" + "\n\n".join(tables_content)
            
            if content.strip():
                print(f"[DEBUG] python-docx成功提取 {len(content)} 字符")
                return content
            else:
                print("[DEBUG] python-docx提取内容为空")
                method_results['python-docx'] = 'Word文档为空或无可读内容'
            
        except Exception as e:
            error_str = str(e)
            print(f"[DEBUG] python-docx提取失败: {error_str}")
            method_results['python-docx'] = error_str
            
            # 如果是.doc文件且python-docx失败，给出详细的错误报告
            if file_path.lower().endswith('.doc'):
                if "not a Word file" in error_str or "themeManager" in error_str:
                    # 生成详细的诊断报告
                    diagnostic_info = []
                    diagnostic_info.append("=== .doc文件处理诊断报告 ===")
                    for method, result in method_results.items():
                        diagnostic_info.append(f"{method}: {result}")
                    diagnostic_info.append(f"重试次数: {retry_count}")
                    diagnostic_info.append(f"python-docx错误: {error_str}")
                    
                    print(f"[DEBUG] 生成诊断报告:")
                    for line in diagnostic_info:
                        print(f"[DEBUG] {line}")
                    
                    if retry_count < max_retries:
                        print(f"[DEBUG] 准备最后重试...")
                        import time
                        time.sleep(2)  # 等待更长时间
                        return self._extract_word_content(file_path, retry_count + 1)
                    else:
                        return f"老版本.doc文件处理失败，详细诊断：\\n" + "\\n".join(diagnostic_info)
            
            return f"Word文档读取失败: {error_str}"
    
    def _extract_doc_content_with_com(self, file_path: str) -> Optional[str]:
        """使用Windows COM接口提取老版本.doc文件内容"""
        if not HAS_WIN32COM:
            raise Exception("win32com.client不可用")
        
        import os
        word_app = None
        document = None
        
        try:
            # 初始化COM库
            import pythoncom
            pythoncom.CoInitialize()
            
            # 获取绝对路径
            abs_file_path = os.path.abspath(file_path)
            if not os.path.exists(abs_file_path):
                return "文件不存在"
            
            # 创建Word应用程序对象
            word_app = win32com.client.Dispatch('Word.Application')
            word_app.Visible = False  # 不显示Word界面
            word_app.DisplayAlerts = False  # 不显示警告对话框
            
            # 打开文档
            document = word_app.Documents.Open(abs_file_path, ReadOnly=True)
            
            # 提取文档内容
            content_parts = []
            
            # 提取正文内容
            if document.Content.Text:
                main_text = document.Content.Text.strip()
                if main_text:
                    content_parts.append("=== 文档内容 (COM接口提取) ===")
                    content_parts.append(main_text)
            
            # 提取表格内容
            if document.Tables.Count > 0:
                content_parts.append("\n\n=== 表格内容 ===")
                
                for i in range(1, document.Tables.Count + 1):  # COM对象索引从1开始
                    table = document.Tables(i)
                    content_parts.append(f"\n[表格 {i}]")
                    
                    try:
                        # 提取表格文本
                        table_text = table.Range.Text.strip()
                        if table_text:
                            # 清理表格文本中的特殊字符
                            table_text = table_text.replace('\r\a', '\n').replace('\a', ' | ')
                            content_parts.append(table_text)
                    except:
                        content_parts.append("表格内容提取失败")
            
            # 提取文档属性
            try:
                built_in_props = document.BuiltInDocumentProperties
                content_parts.append(f"\n\n=== 文档信息 ===")
                
                # 获取一些基本属性
                try:
                    title = built_in_props("Title").Value
                    if title:
                        content_parts.append(f"标题: {title}")
                except:
                    pass
                
                try:
                    author = built_in_props("Author").Value
                    if author:
                        content_parts.append(f"作者: {author}")
                except:
                    pass
                
                try:
                    subject = built_in_props("Subject").Value
                    if subject:
                        content_parts.append(f"主题: {subject}")
                except:
                    pass
                    
            except:
                pass  # 如果无法获取属性，跳过
            
            final_content = "\n".join(content_parts)
            # 清理有问题的Unicode字符
            cleaned_content = self._clean_unicode_text(final_content) if final_content else None
            return cleaned_content if cleaned_content and cleaned_content.strip() else "Word文档为空或无可读内容"
            
        except Exception as e:
            error_msg = str(e)
            print(f"COM接口提取Word文档失败: {error_msg}")
            
            # 提供更具体的错误信息
            if "找不到应用程序" in error_msg or "Application" in error_msg:
                raise Exception("系统未安装Microsoft Word，无法使用COM接口读取.doc文件")
            elif "无法打开" in error_msg or "Open" in error_msg:
                raise Exception("Word无法打开该.doc文件，文件可能损坏或格式不正确")
            else:
                raise Exception(f"COM接口读取失败: {error_msg}")
                
        finally:
            # 确保清理资源
            try:
                if document:
                    document.Close(SaveChanges=False)
                if word_app:
                    word_app.Quit()
            except:
                pass  # 忽略清理过程中的错误
            
            # 清理COM库
            try:
                import pythoncom
                pythoncom.CoUninitialize()
            except:
                pass
    
    def _extract_doc_with_docx2txt(self, file_path: str) -> Optional[str]:
        """使用docx2txt提取文档内容"""
        if not HAS_DOCX2TXT:
            raise Exception("docx2txt不可用")
        
        try:
            content = docx2txt.process(file_path)
            if content and content.strip():
                return f"=== 文档内容 (docx2txt提取) ===\n\n{content.strip()}"
            else:
                return None
        except Exception as e:
            raise Exception(f"docx2txt提取失败: {str(e)}")
    
    def _extract_doc_with_olefile(self, file_path: str) -> Optional[str]:
        """使用olefile直接解析OLE文档结构"""
        if not HAS_OLEFILE:
            raise Exception("olefile不可用")
        
        try:
            if not olefile.isOleFile(file_path):
                raise Exception("文件不是OLE复合文档格式")
            
            ole = olefile.OleFileIO(file_path)
            content_parts = []
            content_parts.append("=== 文档内容 (olefile直接解析) ===\n")
            
            try:
                # 获取文档摘要信息
                if ole.exists('\x05SummaryInformation'):
                    content_parts.append("[文档属性]")
                    # 尝试提取一些基本信息
                    content_parts.append("检测到文档摘要信息流")
                
                # 尝试从WordDocument流提取文本
                if ole.exists('WordDocument'):
                    content_parts.append("\n[主要内容]")
                    
                    # 使用olefile的内部方法读取流数据
                    try:
                        # 方法1：尝试使用_olestream_read (使用list格式的stream名称)
                        word_data = None
                        try:
                            if hasattr(ole, '_olestream_read'):
                                word_data = ole._olestream_read(['WordDocument'], 0, 8192)  # 读取8KB
                        except:
                            try:
                                # 备选方法：直接使用字符串
                                word_data = ole._olestream_read('WordDocument', 0, 8192)
                            except:
                                pass
                            if word_data:
                                # 尝试多种编码解析
                                text_found = False
                                
                                # UTF-16LE解码（Word常用）
                                try:
                                    if len(word_data) > 100:  # 跳过文件头
                                        text_data = word_data[100:]  # 跳过前100字节的头部信息
                                        text = text_data.decode('utf-16le', errors='ignore')
                                        # 清理并提取有意义的文本
                                        cleaned_text = self._clean_extracted_text(text)
                                        if cleaned_text and len(cleaned_text) > 20:
                                            content_parts.append(f"提取的文本内容：\n{cleaned_text}")
                                            text_found = True
                                except:
                                    pass
                                
                                # 如果UTF-16LE失败，尝试其他方法
                                if not text_found:
                                    # 尝试查找ASCII文本模式
                                    ascii_text = ""
                                    for i in range(0, len(word_data) - 1, 2):
                                        char_code = word_data[i]
                                        next_char = word_data[i + 1] if i + 1 < len(word_data) else 0
                                        
                                        # Word文档中，文本通常以UTF-16LE存储（低字节在前）
                                        if 32 <= char_code <= 126 and next_char == 0:  # ASCII字符 + 零字节
                                            ascii_text += chr(char_code)
                                        elif len(ascii_text) > 0 and char_code < 32:
                                            ascii_text += " "  # 用空格替换控制字符
                                    
                                    cleaned_ascii = self._clean_extracted_text(ascii_text)
                                    if cleaned_ascii and len(cleaned_ascii) > 20:
                                        content_parts.append(f"ASCII提取的文本内容：\n{cleaned_ascii}")
                                        text_found = True
                                
                                if not text_found:
                                    content_parts.append("WordDocument流存在，但无法解析可读文本")
                    except Exception as e:
                        content_parts.append(f"WordDocument流解析错误: {str(e)}")
                
                # 尝试从1Table流提取表格数据
                if ole.exists('1Table'):
                    content_parts.append("\n[表格数据]")
                    content_parts.append("检测到表格数据流")
                
                final_content = "\n".join(content_parts)
                return final_content if len(final_content) > 100 else None
                
            finally:
                ole.close()
                
        except Exception as e:
            raise Exception(f"olefile解析失败: {str(e)}")
    
    def _clean_extracted_text(self, raw_text: str) -> str:
        """清理提取的原始文本"""
        if not raw_text:
            return ""
        
        # 移除过多的空白字符
        import re
        
        # 替换多个连续空格为单个空格
        cleaned = re.sub(r'\s+', ' ', raw_text)
        
        # 移除非打印字符（除了常见的空白字符）
        cleaned = ''.join(char for char in cleaned if char.isprintable() or char in '\n\r\t ')
        
        # 移除过短的"词汇"（可能是编码错误）
        words = cleaned.split()
        meaningful_words = []
        for word in words:
            # 保留包含中文字符的词汇，或长度>=2的英文词汇
            if any('\u4e00' <= char <= '\u9fff' for char in word) or len(word) >= 2:
                meaningful_words.append(word)
        
        result = ' '.join(meaningful_words)
        return result.strip()
    
    def _extract_doc_with_system_tools(self, file_path: str) -> Optional[str]:
        """简化版：直接使用LibreOffice处理.doc文件，无健康检查"""
        
        # 直接尝试LibreOffice转换（快速简单）
        result = self._extract_with_libreoffice_simple(file_path)
        if result:
            return result + "\n\n--- LibreOffice提取 ---"
        
        # LibreOffice失败时，尝试一次重试（处理临时问题）
        result = self._extract_with_libreoffice_simple(file_path)
        if result:
            return result + "\n\n--- LibreOffice提取 ---"
        
        # LibreOffice完全失败时尝试其他工具
        return self._extract_with_other_tools(file_path)
    
    def _read_libreoffice_output(self, txt_file: str) -> Optional[str]:
        """读取LibreOffice输出文件并进行优化处理"""
        import os
        
        try:
            # 尝试多种编码读取文件，优先选择最佳内容
            content = None
            best_content = None
            best_score = 0
            
            encodings = ['utf-8', 'gbk', 'utf-16', 'latin1']
            
            for encoding in encodings:
                try:
                    with open(txt_file, 'r', encoding=encoding, errors='ignore') as f:
                        test_content = f.read().strip()
                        if len(test_content) > 10:
                            # 计算内容质量分数
                            chinese_chars = sum(1 for c in test_content if 0x4e00 <= ord(c) <= 0x9fff)
                            ascii_chars = sum(1 for c in test_content if ord(c) < 128)
                            total_chars = len(test_content)
                            
                            # 质量评分：中文字符数 + ASCII可读字符数
                            score = chinese_chars * 2 + ascii_chars * 0.5
                            
                            print(f"[{encoding}] 长度:{total_chars}, 中文:{chinese_chars}, ASCII:{ascii_chars}, 分数:{score}")
                            
                            if score > best_score:
                                best_content = test_content
                                best_score = score
                                print(f"[{encoding}] 当前最佳选择")
                except Exception as e:
                    print(f"[{encoding}] 读取失败: {e}")
                    continue
            
            content = best_content
            
            # 清理临时文件
            try:
                os.remove(txt_file)
            except:
                pass
            
            if content:
                # 清理有问题的Unicode字符，确保兼容性
                cleaned_content = self._clean_unicode_text(content)
                final_result = f"=== 文档内容 (LibreOffice提取) ===\n\n{cleaned_content}"
                print(f"[OK] LibreOffice内容提取完成，最终长度: {len(final_result)}")
                return final_result
            else:
                print("[FAIL] LibreOffice输出文件读取失败或内容为空")
                return None
                
        except Exception as e:
            print(f"[ERROR] 读取LibreOffice输出文件异常: {str(e)}")
            return None
    
    def _find_libreoffice_path(self) -> Optional[str]:
        """快速查找LibreOffice可执行文件路径"""
        import subprocess
        import os
        
        # 常见的LibreOffice路径，按优先级排序
        possible_paths = [
            r'C:\Program Files\LibreOffice\program\soffice.exe',      # Windows 64位
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe', # Windows 32位
            '/usr/bin/libreoffice',   # Linux 标准路径
            '/usr/bin/soffice',       # Linux 备用路径
            'soffice'                 # PATH环境变量
        ]
        
        # 快速检查，不做复杂验证
        for path in possible_paths:
            try:
                if path.startswith('/') or path.startswith('C:'):
                    # 绝对路径检查
                    if os.path.exists(path):
                        return path
                else:
                    # PATH环境变量检查
                    try:
                        subprocess.run([path, '--version'], capture_output=True, timeout=3)
                        return path
                    except (subprocess.TimeoutExpired, FileNotFoundError):
                        continue
            except:
                continue
        
        return None
    
    def _extract_with_libreoffice_simple(self, file_path: str) -> Optional[str]:
        """简化版LibreOffice文档转换，无健康检查"""
        import subprocess
        import os
        
        # 快速找到LibreOffice路径
        libreoffice_path = self._find_libreoffice_path()
        if not libreoffice_path:
            return None
        
        try:
            # 创建临时输出目录
            temp_dir = os.path.join(os.path.dirname(file_path), 'temp')
            os.makedirs(temp_dir, exist_ok=True)
            
            # 使用优化的后台运行参数
            convert_cmd = [
                libreoffice_path,
                '--headless',              # 无界面模式
                '--invisible',             # 不显示启动logo和程序窗口
                '--nologo',                # 禁用启动画面
                '--nofirststartwizard',    # 禁用首次启动向导
                '--norestore',             # 禁止崩溃后自动恢复
                '--nolockcheck',           # 禁用文件锁定检查
                '--nodefault',             # 不创建默认文档
                '--convert-to', 'txt',     # 转换为文本格式
                '--outdir', temp_dir,      # 输出目录
                os.path.abspath(file_path) # 使用绝对路径
            ]
            
            # 执行转换，设置合理超时
            result = subprocess.run(convert_cmd, capture_output=True, text=True, timeout=30)
            
            if result.returncode == 0:
                # 查找生成的输出文件
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                txt_file = os.path.join(temp_dir, f"{base_name}.txt")
                
                if os.path.exists(txt_file):
                    # 读取并处理输出文件
                    content = self._read_libreoffice_output(txt_file)
                    if content:
                        return content
            
            return None
            
        except subprocess.TimeoutExpired:
            return None
        except Exception as e:
            return None
    
    def _extract_with_other_tools(self, file_path: str) -> Optional[str]:
        """使用其他系统工具提取文档内容"""
        import subprocess
        
        # 尝试其他工具
        other_tools = [
            # Pandoc转换
            {
                'name': 'Pandoc',
                'command': ['pandoc', '-f', 'doc', '-t', 'plain', file_path],
                'check_cmd': ['pandoc', '--version']
            }
        ]
        
        for tool in other_tools:
            try:
                # 检查工具是否可用
                print(f"检查 {tool['name']} 可用性...")
                result = subprocess.run(tool['check_cmd'], capture_output=True, text=True, timeout=10)
                if result.returncode == 0:
                    print(f"[INFO] {tool['name']} 可用，尝试处理...")
                    
                    # 执行转换命令
                    conv_result = subprocess.run(tool['command'], capture_output=True, text=True, timeout=30)
                    if conv_result.returncode == 0 and conv_result.stdout.strip():
                        # 清理有问题的Unicode字符
                        cleaned_content = self._clean_unicode_text(conv_result.stdout.strip())
                        final_result = f"=== 文档内容 ({tool['name']}提取) ===\n\n{cleaned_content}"
                        print(f"[OK] {tool['name']} 处理成功")
                        return final_result
                    else:
                        print(f"[FAIL] {tool['name']} 处理失败")
                else:
                    print(f"[SKIP] {tool['name']} 不可用")
                    
            except (subprocess.TimeoutExpired, subprocess.CalledProcessError, FileNotFoundError):
                print(f"[FAIL] {tool['name']} 执行失败")
                continue
            except Exception as e:
                print(f"[ERROR] {tool['name']} 执行异常: {str(e)}")
                continue
        
        print("[FAIL] 所有系统工具都无法处理该文档")
        return None
    
    def _clean_unicode_text(self, text: str) -> str:
        """清理文本中有问题的Unicode字符，确保编码兼容性"""
        if not text:
            return ""
        
        # 直接使用最严格的字符过滤方法
        import re
        
        try:
            # 方法1：测试当前字符串是否可以安全编码
            try:
                text.encode('utf-8')
                text.encode('gbk', errors='strict')
                # 如果都成功，直接返回
                return text
            except UnicodeEncodeError:
                pass  # 继续进行字符清理
            
            # 方法2：字符级别的严格过滤 - 只保留绝对安全的字符
            safe_chars = []
            for char in text:
                char_code = ord(char)
                
                # 只保留ASCII字符（绝对安全）
                if char_code < 128:
                    safe_chars.append(char)
                # 只保留最核心的中文字符范围（基本汉字区）
                elif 0x4e00 <= char_code <= 0x9fff:
                    # 直接测试GBK兼容性
                    try:
                        char.encode('gbk', errors='strict')
                        safe_chars.append(char)
                    except UnicodeEncodeError:
                        safe_chars.append('?')  # 替换不兼容的中文字符
                # 只保留最常用的中文标点（经过GBK测试）
                elif char in '，。！？；：':
                    try:
                        char.encode('gbk', errors='strict')
                        safe_chars.append(char)
                    except UnicodeEncodeError:
                        if char == '，':
                            safe_chars.append(',')
                        elif char == '。':
                            safe_chars.append('.')
                        elif char == '！':
                            safe_chars.append('!')
                        elif char == '？':
                            safe_chars.append('?')
                        elif char == '；':
                            safe_chars.append(';')
                        elif char == '：':
                            safe_chars.append(':')
                        else:
                            safe_chars.append('.')
                # 保留基本空白字符
                elif char in '\n\r\t ':
                    safe_chars.append(char)
                # 所有其他字符都过滤掉或替换
                else:
                    # 已知问题字符直接替换
                    if char_code == 0x023b:  # Ȼ
                        safe_chars.append('c')  
                    elif char_code == 0x2717:  # [FAIL]
                        safe_chars.append('x')
                    elif char_code == 0x2713:  # [OK]
                        safe_chars.append('v')
                    elif char_code == 0x2717:  # [FAIL]
                        safe_chars.append('x')
                    elif 0x2000 <= char_code <= 0x206f:  # 通用标点符号
                        safe_chars.append(' ')
                    elif 0x2700 <= char_code <= 0x27bf:  # 符号和装饰符号
                        safe_chars.append('*')
                    elif 0x0100 <= char_code <= 0x017f:  # 拉丁文扩展字符
                        safe_chars.append('?')
                    else:
                        safe_chars.append(' ')  # 其他特殊字符替换为空格
            
            result = ''.join(safe_chars)
            
            # 清理多余的空格和替换字符
            result = re.sub(r'\s+', ' ', result)  # 多个空格合并为一个
            result = re.sub(r'\?{2,}', '?', result)  # 多个?合并
            result = re.sub(r'\?\s+\?', '?', result)  # 清理?之间的空格
            
            # 最终测试
            try:
                result.encode('utf-8')
                result.encode('gbk', errors='strict')
                return result.strip()
            except UnicodeEncodeError:
                # 如果还有问题，使用最保守的方法
                ascii_result = result.encode('ascii', errors='ignore').decode('ascii')
                return ascii_result.strip()
                
        except Exception as e:
            print(f"Unicode清理失败，使用最保守方法: {e}")
            # 最保守方法：只保留基本ASCII字符
            try:
                return text.encode('ascii', errors='ignore').decode('ascii').strip()
            except:
                return "文档内容包含特殊字符，无法正常显示"
    
    def _extract_excel_content(self, file_path: str) -> Optional[str]:
        """提取Excel文件内容"""
        if not pd:
            return "需要安装pandas库才能读取Excel文件"
        
        try:
            # 尝试多种引擎读取Excel文件
            engines = ['openpyxl', 'xlrd', None]  # openpyxl用于.xlsx，xlrd用于.xls
            
            for engine in engines:
                try:
                    # 读取所有工作表
                    if engine:
                        excel_file = pd.ExcelFile(file_path, engine=engine)
                    else:
                        excel_file = pd.ExcelFile(file_path)
                    
                    sheets_content = []
                    
                    for sheet_name in excel_file.sheet_names:
                        try:
                            if engine:
                                df = pd.read_excel(file_path, sheet_name=sheet_name, engine=engine)
                            else:
                                df = pd.read_excel(file_path, sheet_name=sheet_name)
                            
                            # 转换为字符串表示
                            sheet_text = f"[工作表: {sheet_name}]\n"
                            sheet_text += df.to_string(index=False, na_rep='')
                            sheets_content.append(sheet_text)
                            
                        except Exception as e:
                            sheets_content.append(f"[工作表: {sheet_name}] 读取失败: {str(e)}")
                    
                    if sheets_content:
                        return "\n\n".join(sheets_content)
                    
                except Exception as e:
                    print(f"使用引擎 {engine} 读取Excel失败: {e}")
                    continue
            
            return "Excel文件读取失败: 尝试了所有可用的引擎"
            
        except Exception as e:
            print(f"提取Excel文件失败: {e}")
            return f"Excel文件读取失败: {str(e)}"
    
    def _extract_powerpoint_content(self, file_path: str) -> Optional[str]:
        """提取PowerPoint文件内容"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            slides_content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"[幻灯片 {slide_num}]\n"
                
                # 提取文本框内容
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                
                if texts:
                    slide_text += "\n".join(texts)
                else:
                    slide_text += "(无文本内容)"
                
                slides_content.append(slide_text)
            
            return "\n\n".join(slides_content)
            
        except ImportError:
            return "需要安装python-pptx库才能读取PowerPoint文件"
        except Exception as e:
            print(f"提取PowerPoint文件失败: {e}")
            return f"PowerPoint文件读取失败: {str(e)}"
    
    def extract_file_content_partial(self, file_path: str, max_chars: int = 10000) -> Optional[str]:
        """提取文件内容的前N个字符，用于大文件快速搜索"""
        try:
            if not os.path.exists(file_path):
                return None
            
            file_ext = Path(file_path).suffix.lower()
            
            # 对于文本文件，直接读取前N个字符 - 使用智能编码检测
            if file_ext in {'.txt', '.py', '.js', '.html', '.xml', '.yml', '.yaml', '.csv', '.rtf', 
                           '.conf', '.config', '.cfg', '.ini', '.properties', '.env', '.md'}:
                # 首先检测文件编码
                detected_encoding, confidence = EncodingDetector.detect_encoding(file_path)
                try:
                    with open(file_path, 'r', encoding=detected_encoding) as f:
                        content = f.read(max_chars)
                except UnicodeDecodeError:
                    # 如果检测的编码失败，使用UTF-8忽略错误
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read(max_chars)
                    if len(content) == max_chars:
                        # 截断到最后一个完整的行
                        last_newline = content.rfind('\n')
                        if last_newline > 0:
                            content = content[:last_newline]
                    return content
            
            # JSON文件部分读取
            elif file_ext == '.json':
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read(max_chars)
                    return content
            
            # 对于复杂文件类型，跳过部分读取，返回None
            else:
                return None
                
        except Exception as e:
            print(f"部分提取文件内容失败 {file_path}: {str(e)}")
            return None
    
    def highlight_text(self, text: str, keyword: str) -> str:
        """在文本中高亮关键词"""
        try:
            pattern = re.compile(re.escape(keyword), re.IGNORECASE)
            return pattern.sub(lambda m: f"<mark>{m.group()}</mark>", text)
        except Exception:
            return text
    
    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """获取文件信息"""
        try:
            if not os.path.exists(file_path):
                return {}
            
            stat = os.stat(file_path)
            return {
                "size": stat.st_size,
                "modified_time": stat.st_mtime,
                "extension": Path(file_path).suffix.lower(),
                "name": Path(file_path).name,
                "is_searchable": Path(file_path).suffix.lower() in self.supported_extensions
            }
        except Exception:
            return {}
    
    def search_multiple_files(self, file_paths: List[str], keyword: str) -> Dict[str, List[Dict[str, Any]]]:
        """在多个文件中搜索"""
        results = {}
        
        for file_path in file_paths:
            matches = self.search_in_file(file_path, keyword)
            if matches:
                results[file_path] = matches
        
        return results
    
    def fuzzy_search(self, text: str, keyword: str, threshold: float = 0.6) -> List[Dict[str, Any]]:
        """模糊搜索"""
        matches = []
        lines = text.split('\n')
        
        for line_num, line in enumerate(lines, 1):
            # 简单的模糊匹配算法
            similarity = self._calculate_similarity(line.lower(), keyword.lower())
            
            if similarity >= threshold:
                matches.append({
                    "line_number": line_num,
                    "text": line,
                    "similarity": similarity
                })
        
        return sorted(matches, key=lambda x: x["similarity"], reverse=True)
    
    def _calculate_similarity(self, str1: str, str2: str) -> float:
        """计算字符串相似度"""
        if not str1 or not str2:
            return 0.0
        
        # 简单的包含匹配
        if str2 in str1:
            return len(str2) / len(str1)
        
        # 计算公共子序列长度
        common_chars = set(str1) & set(str2)
        if not common_chars:
            return 0.0
        
        return len(common_chars) / max(len(set(str1)), len(set(str2)))
    
    def _extract_office_theme_content(self, file_path: str) -> Optional[str]:
        """提取Office主题文件的内容"""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            content_parts = []
            content_parts.append("=== Office主题文件内容 ===\n")
            
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # 提取主题信息
                if 'theme/theme/theme1.xml' in zip_file.namelist():
                    with zip_file.open('theme/theme/theme1.xml') as f:
                        xml_content = f.read().decode('utf-8', errors='replace')
                        
                        try:
                            root = ET.fromstring(xml_content)
                            
                            # 提取主题名称
                            theme_name = root.get('name', '未命名主题')
                            # 清理名称中的特殊字符
                            theme_name = theme_name.replace('\u200b', '').replace('\u200c', '').replace('\u200d', '')
                            content_parts.append(f"主题名称: {theme_name}\n")
                            
                            # 提取颜色方案
                            color_schemes = []
                            for clr_scheme in root.iter():
                                if 'clrScheme' in clr_scheme.tag:
                                    scheme_name = clr_scheme.get('name', '默认')
                                    color_schemes.append(f"颜色方案: {scheme_name}")
                                    
                                    # 提取具体颜色
                                    colors = []
                                    for color_elem in clr_scheme:
                                        color_name = color_elem.tag.split('}')[-1] if '}' in color_elem.tag else color_elem.tag
                                        
                                        # 查找颜色值
                                        for child in color_elem:
                                            if 'srgbClr' in child.tag:
                                                color_val = child.get('val', '')
                                                if color_val:
                                                    colors.append(f"  {color_name}: #{color_val}")
                                                    break
                                            elif 'sysClr' in child.tag:
                                                sys_color = child.get('val', '')
                                                last_color = child.get('lastClr', '')
                                                if last_color:
                                                    colors.append(f"  {color_name}: {sys_color} (#{last_color})")
                                                else:
                                                    colors.append(f"  {color_name}: {sys_color}")
                                                break
                                    
                                    if colors:
                                        color_schemes.extend(colors)
                            
                            if color_schemes:
                                content_parts.append("\n颜色配置:")
                                content_parts.extend([f"\n{scheme}" for scheme in color_schemes])
                            
                            # 提取字体方案
                            font_schemes = []
                            for font_scheme in root.iter():
                                if 'fontScheme' in font_scheme.tag:
                                    scheme_name = font_scheme.get('name', '默认')
                                    font_schemes.append(f"字体方案: {scheme_name}")
                                    
                                    # 提取字体信息
                                    fonts = []
                                    for font_elem in font_scheme.iter():
                                        if 'font' in font_elem.tag:
                                            script = font_elem.get('script', '')
                                            typeface = font_elem.get('typeface', '')
                                            if typeface:
                                                # 清理字体名称中的特殊字符
                                                typeface = typeface.encode('ascii', errors='ignore').decode('ascii')
                                                if typeface:  # 确保清理后不为空
                                                    if script:
                                                        fonts.append(f"  {script}: {typeface}")
                                                    else:
                                                        fonts.append(f"  字体: {typeface}")
                                    
                                    if fonts:
                                        font_schemes.extend(fonts)
                            
                            if font_schemes:
                                content_parts.append("\n\n字体配置:")
                                content_parts.extend([f"\n{font}" for font in font_schemes])
                            
                            # 提取格式方案
                            format_schemes = []
                            for fmt_scheme in root.iter():
                                if 'fmtScheme' in fmt_scheme.tag:
                                    scheme_name = fmt_scheme.get('name', '默认')
                                    format_schemes.append(f"格式方案: {scheme_name}")
                            
                            if format_schemes:
                                content_parts.append("\n\n格式配置:")
                                content_parts.extend([f"\n{fmt}" for fmt in format_schemes])
                            
                        except ET.ParseError as e:
                            content_parts.append(f"\nXML解析错误: {e}")
                        except Exception as e:
                            content_parts.append(f"\n主题解析错误: {e}")
                
                # 提取内容类型信息
                if '[Content_Types].xml' in zip_file.namelist():
                    content_parts.append("\n\n文件类型信息:")
                    try:
                        with zip_file.open('[Content_Types].xml') as f:
                            xml_content = f.read().decode('utf-8', errors='replace')
                            root = ET.fromstring(xml_content)
                            
                            for override in root.iter():
                                if 'Override' in override.tag:
                                    part_name = override.get('PartName', '')
                                    content_type = override.get('ContentType', '')
                                    if part_name and content_type:
                                        content_parts.append(f"\n  {part_name}: {content_type.split('.')[-1]}")
                    except:
                        pass
            
            # 添加说明
            content_parts.append(f"\n\n文件类型: Microsoft Office主题文件")
            content_parts.append(f"说明: 此文件包含Office应用程序的主题配置信息，包括颜色、字体和格式方案。")
            
            return "".join(content_parts)
            
        except Exception as e:
            print(f"提取Office主题文件失败: {e}")
            return None